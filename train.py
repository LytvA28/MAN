from pptx import Presentation
import csv
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.linear_model import LogisticRegression
from sklearn.pipeline import Pipeline
from sklearn.model_selection import train_test_split
from sklearn.metrics import classification_report
import joblib
import os

def extract_text_from_pptx(path):
    prs = Presentation(path)
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                if shape.text.strip():
                    texts.append(shape.text.strip())
    return "\n".join(texts)

def load_dataset(csv_path):
    texts = []
    labels = []
    with open(csv_path, encoding="utf-8") as f:
        reader = csv.DictReader(f)
        for row in reader:
            text = extract_text_from_pptx(row["path"])
            texts.append(text)
            labels.append(row["label"])
    return texts, labels

def train_model():
    print("Завантаження датасету...")
    texts, labels = load_dataset("dataset.csv")

    print("Розділення на train/test...")
    X_train, X_test, y_train, y_test = train_test_split(
        texts, labels, test_size=0.2, random_state=42, stratify=labels
    )

    print("Навчання моделі...")
    model = Pipeline([
        ("tfidf", TfidfVectorizer(max_features=20000, ngram_range=(1, 2))),
        ("clf", LogisticRegression(max_iter=300))
    ])

    model.fit(X_train, y_train)

    print("Оцінка моделі:")
    y_pred = model.predict(X_test)
    print(classification_report(y_test, y_pred))

    joblib.dump(model, "ai_presentation_detector.joblib")
    print("Модель збережена як ai_presentation_detector.joblib")

if __name__ == "__main__":
    train_model()
